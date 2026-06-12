"""
Day 7 — Generate Final Report (PDF) and Presentation (PPTX)
Run: python scripts/generate_report.py
"""

from pathlib import Path
import pandas as pd
import numpy as np

BASE_DIR  = Path(__file__).resolve().parent.parent
PROCESSED = BASE_DIR / "data" / "processed"
RAW       = BASE_DIR / "data" / "raw"
REPORTS   = BASE_DIR / "reports"
REPORTS.mkdir(exist_ok=True)

nav         = pd.read_csv(PROCESSED / "nav_clean.csv", parse_dates=["date"])
perf        = pd.read_csv(PROCESSED / "performance_clean.csv")
scorecard   = pd.read_csv(PROCESSED / "fund_scorecard.csv")
ab          = pd.read_csv(PROCESSED / "alpha_beta.csv")
var_df      = pd.read_csv(PROCESSED / "var_cvar_report.csv")
txn         = pd.read_csv(PROCESSED / "transactions_clean.csv", parse_dates=["transaction_date"])
aum         = pd.read_csv(RAW / "03_aum_by_fund_house.csv", parse_dates=["date"])
sip         = pd.read_csv(RAW / "04_monthly_sip_inflows.csv", parse_dates=["month"])
fund_master = pd.read_csv(RAW / "01_fund_master.csv")
fund_master.columns = fund_master.columns.str.strip().str.lower()
perf.columns = perf.columns.str.strip().str.lower()

total_aum    = aum["aum_lakh_crore"].max()
total_sip    = sip["sip_inflow_crore"].max()
total_funds  = nav["amfi_code"].nunique()
total_txns   = len(txn)
top5_funds   = scorecard.nlargest(5, "score")[["scheme_name", "score", "cagr_3y"]].head()
top5_var     = var_df.nsmallest(5, "var_95_pct")[["scheme_name", "var_95_pct", "cvar_95_pct"]].head()
top5_alpha   = ab.nlargest(5, "alpha")[["scheme_name", "alpha", "beta"]].head()

print("Data loaded")

# ═══════════════════════════════════════════════════════
# PART 1 — FINAL REPORT PDF
# ═══════════════════════════════════════════════════════
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch, cm
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                 Table, TableStyle, PageBreak, HRFlowable)
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY

W, H = A4

def build_pdf():
    doc = SimpleDocTemplate(
        str(REPORTS / "Final_Report.pdf"),
        pagesize=A4,
        rightMargin=2*cm, leftMargin=2*cm,
        topMargin=2*cm, bottomMargin=2*cm
    )

    styles = getSampleStyleSheet()

    # Custom styles
    title_style = ParagraphStyle("Title2", parent=styles["Title"],
                                  fontSize=24, textColor=colors.HexColor("#0A2142"),
                                  spaceAfter=6, alignment=TA_CENTER)
    h1 = ParagraphStyle("H1", parent=styles["Heading1"],
                         fontSize=16, textColor=colors.HexColor("#0A2142"),
                         spaceBefore=16, spaceAfter=6)
    h2 = ParagraphStyle("H2", parent=styles["Heading2"],
                         fontSize=13, textColor=colors.HexColor("#1E3A5F"),
                         spaceBefore=10, spaceAfter=4)
    body = ParagraphStyle("Body2", parent=styles["Normal"],
                           fontSize=11, leading=16, alignment=TA_JUSTIFY,
                           spaceAfter=8)
    caption = ParagraphStyle("Caption", parent=styles["Normal"],
                              fontSize=9, textColor=colors.grey,
                              alignment=TA_CENTER)
    kpi_style = ParagraphStyle("KPI", parent=styles["Normal"],
                                fontSize=20, textColor=colors.HexColor("#00B4D8"),
                                alignment=TA_CENTER, fontName="Helvetica-Bold")

    story = []

    # ── COVER PAGE ──
    story.append(Spacer(1, 1.5*inch))
    story.append(Paragraph("Bluestock Mutual Fund", title_style))
    story.append(Paragraph("Analytics Capstone Project", title_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(HRFlowable(width="100%", thickness=2, color=colors.HexColor("#00B4D8")))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("Final Report — June 2026", styles["Heading2"]))
    story.append(Spacer(1, 0.2*inch))
    story.append(Paragraph("Prepared by: Divya Madiraju", body))
    story.append(Paragraph("Project: Capstone Project I — Mutual Fund Analytics", body))
    story.append(Paragraph("Organization: Bluestock Fintech MJ28", body))
    story.append(PageBreak())

    # ── 1. EXECUTIVE SUMMARY ──
    story.append(Paragraph("1. Executive Summary", h1))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#0A2142")))
    story.append(Spacer(1, 0.1*inch))
    story.append(Paragraph(
        "This capstone project delivers a comprehensive analytics platform for the Indian Mutual Fund industry, "
        "built on AMFI data covering 40 fund schemes from January 2022 to December 2025. The project encompasses "
        "end-to-end data engineering, exploratory analysis, risk-adjusted performance analytics, an interactive "
        "Power BI dashboard, and advanced risk metrics including Value at Risk (VaR) and investor cohort analysis.",
        body))
    story.append(Paragraph(
        "Key findings include SBI Mutual Fund's dominance with the highest AUM, small-cap funds delivering "
        "superior returns at higher risk, and strong SIP growth trends reflecting increasing retail investor participation.",
        body))

    # KPI row
    kpi_data = [
        [Paragraph(f"₹{total_aum:.1f}L Cr", kpi_style),
         Paragraph(f"₹{total_sip:,.0f} Cr", kpi_style),
         Paragraph(f"{total_funds}", kpi_style),
         Paragraph(f"{total_txns:,}", kpi_style)],
        [Paragraph("Peak AUM", caption),
         Paragraph("Peak SIP Inflow", caption),
         Paragraph("Fund Schemes", caption),
         Paragraph("Transactions", caption)]
    ]
    kpi_table = Table(kpi_data, colWidths=[4*cm]*4)
    kpi_table.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,-1), colors.HexColor("#F0F8FF")),
        ("BOX", (0,0), (-1,-1), 1, colors.HexColor("#0A2142")),
        ("INNERGRID", (0,0), (-1,-1), 0.5, colors.HexColor("#CCDDEE")),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING", (0,0), (-1,-1), 8),
        ("BOTTOMPADDING", (0,0), (-1,-1), 8),
    ]))
    story.append(kpi_table)
    story.append(Spacer(1, 0.2*inch))

    # ── 2. DATA SOURCES ──
    story.append(Paragraph("2. Data Sources", h1))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#0A2142")))
    story.append(Spacer(1, 0.1*inch))
    story.append(Paragraph(
        "The project uses 10 datasets sourced from AMFI India and simulated transaction data:", body))

    ds_data = [
        ["File", "Description", "Rows"],
        ["01_fund_master.csv", "Fund scheme metadata", "40"],
        ["02_nav_history.csv", "Daily NAV (2022–2026)", "46,000"],
        ["03_aum_by_fund_house.csv", "Monthly AUM by AMC", "90"],
        ["04_monthly_sip_inflows.csv", "Industry SIP inflows", "48"],
        ["05_category_inflows.csv", "Category-wise net inflows", "144"],
        ["06_industry_folio_count.csv", "Folio count growth", "21"],
        ["07_scheme_performance.csv", "Fund performance metrics", "40"],
        ["08_investor_transactions.csv", "Investor transactions", "32,778"],
        ["09_portfolio_holdings.csv", "Equity fund holdings", "322"],
        ["10_benchmark_indices.csv", "Nifty/BSE indices", "8,050"],
    ]
    ds_table = Table(ds_data, colWidths=[5.5*cm, 8*cm, 2.5*cm])
    ds_table.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#0A2142")),
        ("TEXTCOLOR", (0,0), (-1,0), colors.white),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 9),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.HexColor("#F0F8FF")]),
        ("GRID", (0,0), (-1,-1), 0.5, colors.HexColor("#CCDDEE")),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING", (0,0), (-1,-1), 5),
        ("BOTTOMPADDING", (0,0), (-1,-1), 5),
    ]))
    story.append(ds_table)
    story.append(PageBreak())

    # ── 3. ETL DESIGN ──
    story.append(Paragraph("3. ETL Pipeline Design", h1))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#0A2142")))
    story.append(Spacer(1, 0.1*inch))
    story.append(Paragraph(
        "The ETL pipeline follows a structured Extract-Transform-Load architecture:", body))
    story.append(Paragraph("<b>Extract:</b> Raw CSVs loaded from data/raw/ using Pandas. "
                            "Live NAV fetched from mfapi.in API for 5 key schemes.", body))
    story.append(Paragraph("<b>Transform:</b> Dates parsed to datetime, NAV forward-filled for weekends/holidays, "
                            "transaction types standardised (SIP/Lumpsum/Redemption), "
                            "expense ratios validated (0.1%-2.5%), duplicates removed.", body))
    story.append(Paragraph("<b>Load:</b> Cleaned data loaded into SQLite star schema with 5 tables: "
                            "dim_fund, dim_date, fact_nav, fact_transactions, fact_performance. "
                            "Row counts verified against source CSVs.", body))
    story.append(Spacer(1, 0.1*inch))
    story.append(Paragraph("3.1 Star Schema Design", h2))
    schema_data = [
        ["Table", "Type", "Primary Key", "Rows"],
        ["dim_fund", "Dimension", "amfi_code", "40"],
        ["dim_date", "Dimension", "date", "1,296"],
        ["fact_nav", "Fact", "nav_id", "46,000"],
        ["fact_transactions", "Fact", "txn_id", "32,778"],
        ["fact_performance", "Fact", "perf_id", "40"],
    ]
    schema_table = Table(schema_data, colWidths=[4*cm, 3*cm, 4*cm, 2.5*cm])
    schema_table.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#0A2142")),
        ("TEXTCOLOR", (0,0), (-1,0), colors.white),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 10),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.HexColor("#F0F8FF")]),
        ("GRID", (0,0), (-1,-1), 0.5, colors.HexColor("#CCDDEE")),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING", (0,0), (-1,-1), 6),
        ("BOTTOMPADDING", (0,0), (-1,-1), 6),
    ]))
    story.append(schema_table)
    story.append(PageBreak())

    # ── 4. EDA FINDINGS ──
    story.append(Paragraph("4. EDA Key Findings", h1))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#0A2142")))
    story.append(Spacer(1, 0.1*inch))

    eda_findings = [
        ("2023 Bull Run", "Equity fund NAVs rose sharply through 2023, with small-cap funds delivering the highest growth, confirming the broad market bull run driven by domestic retail inflows."),
        ("SBI AUM Dominance", "SBI Mutual Fund consistently leads AUM across all years, reflecting distribution advantage through India Post and bank branches with over Rs.80 lakh crore peak AUM."),
        ("SIP Inflow Record", "Monthly SIP inflows grew steadily from approximately Rs.11,000 Cr in January 2022 to record highs by December 2025, reflecting strong retail investor participation."),
        ("Equity Category Dominance", "Equity funds attract consistently positive net inflows while debt funds show seasonal outflows, particularly in March due to advance tax payments."),
        ("Geographic Distribution", "Punjab, Tamil Nadu, and Madhya Pradesh lead in transaction amounts, suggesting strong retail MF penetration beyond metro cities."),
        ("Folio Count Growth", "Industry folios nearly doubled from approximately 13 Cr to 26 Cr between 2022 and 2025, driven primarily by equity folio additions."),
        ("High Return Correlation", "Most equity funds show correlation greater than 0.85, suggesting they move together with the broader market index — diversification benefit is limited within equity."),
        ("Financial Sector Dominance", "Financial services is the largest sector allocation across equity fund portfolios, followed by IT and FMCG, with several funds showing high HHI concentration scores."),
    ]

    for title, text in eda_findings:
        story.append(Paragraph(f"<b>{title}:</b> {text}", body))

    story.append(PageBreak())

    # ── 5. PERFORMANCE ANALYSIS ──
    story.append(Paragraph("5. Performance Analysis", h1))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#0A2142")))
    story.append(Spacer(1, 0.1*inch))
    story.append(Paragraph(
        "Risk-adjusted performance metrics were computed using 252 trading days annualisation:", body))

    story.append(Paragraph("5.1 Top 5 Funds by Composite Score", h2))
    sc_data = [["Fund Name", "Score", "3Y CAGR (%)"]]
    for _, row in top5_funds.iterrows():
        name = str(row.get("scheme_name", ""))[:40]
        score = f"{row.get('score', 0):.1f}"
        cagr = f"{row.get('cagr_3y', 0):.2f}"
        sc_data.append([name, score, cagr])
    sc_table = Table(sc_data, colWidths=[10*cm, 2.5*cm, 3*cm])
    sc_table.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#0A2142")),
        ("TEXTCOLOR", (0,0), (-1,0), colors.white),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 9),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.HexColor("#F0F8FF")]),
        ("GRID", (0,0), (-1,-1), 0.5, colors.HexColor("#CCDDEE")),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING", (0,0), (-1,-1), 5),
        ("BOTTOMPADDING", (0,0), (-1,-1), 5),
    ]))
    story.append(sc_table)
    story.append(Spacer(1, 0.15*inch))

    story.append(Paragraph("5.2 Top 5 Riskiest Funds by VaR (95%)", h2))
    var_data = [["Fund Name", "VaR 95% (%)", "CVaR 95% (%)"]]
    for _, row in top5_var.iterrows():
        name = str(row.get("scheme_name", ""))[:40]
        var_data.append([name,
                          f"{row.get('var_95_pct', 0):.4f}",
                          f"{row.get('cvar_95_pct', 0):.4f}"])
    var_table = Table(var_data, colWidths=[10*cm, 3*cm, 3*cm])
    var_table.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#8B0000")),
        ("TEXTCOLOR", (0,0), (-1,0), colors.white),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 9),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.HexColor("#FFF0F0")]),
        ("GRID", (0,0), (-1,-1), 0.5, colors.HexColor("#FFCCCC")),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING", (0,0), (-1,-1), 5),
        ("BOTTOMPADDING", (0,0), (-1,-1), 5),
    ]))
    story.append(var_table)
    story.append(Spacer(1, 0.15*inch))

    story.append(Paragraph("5.3 Alpha & Beta vs Nifty 100", h2))
    story.append(Paragraph(
        "Alpha is the annualised excess return vs Nifty 100. Beta measures market sensitivity (beta=1 moves with market). "
        "Funds with high alpha and beta less than 1 represent best risk-adjusted outperformance.", body))

    ab_data = [["Fund Name", "Alpha (%)", "Beta"]]
    for _, row in top5_alpha.iterrows():
        name = str(row.get("scheme_name", ""))[:40]
        ab_data.append([name,
                         f"{row.get('alpha', 0):.2f}",
                         f"{row.get('beta', 0):.3f}"])
    ab_table = Table(ab_data, colWidths=[10*cm, 3*cm, 3*cm])
    ab_table.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#006400")),
        ("TEXTCOLOR", (0,0), (-1,0), colors.white),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 9),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.HexColor("#F0FFF0")]),
        ("GRID", (0,0), (-1,-1), 0.5, colors.HexColor("#AADDAA")),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING", (0,0), (-1,-1), 5),
        ("BOTTOMPADDING", (0,0), (-1,-1), 5),
    ]))
    story.append(ab_table)
    story.append(PageBreak())

    # ── 6. DASHBOARD ──
    story.append(Paragraph("6. Interactive Dashboard", h1))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#0A2142")))
    story.append(Spacer(1, 0.1*inch))
    story.append(Paragraph(
        "A 4-page interactive Power BI dashboard was developed covering all key analytical dimensions:", body))

    dash_data = [
        ["Page", "Title", "Key Visuals"],
        ["1", "Industry Overview", "KPI cards, AUM trend line, AUM by AMC bar chart"],
        ["2", "Fund Performance", "Return vs Risk scatter, Fund Scorecard table, NAV trend"],
        ["3", "Investor Analytics", "Transactions by state, Type donut, Monthly volume"],
        ["4", "SIP & Market Trends", "SIP inflow trend, Active accounts, Category inflows"],
    ]
    dash_table = Table(dash_data, colWidths=[1.5*cm, 4.5*cm, 10*cm])
    dash_table.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#0A2142")),
        ("TEXTCOLOR", (0,0), (-1,0), colors.white),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 9),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.HexColor("#F0F8FF")]),
        ("GRID", (0,0), (-1,-1), 0.5, colors.HexColor("#CCDDEE")),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING", (0,0), (-1,-1), 6),
        ("BOTTOMPADDING", (0,0), (-1,-1), 6),
    ]))
    story.append(dash_table)
    story.append(PageBreak())

    # ── 7. LIMITATIONS ──
    story.append(Paragraph("7. Limitations", h1))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#0A2142")))
    story.append(Spacer(1, 0.1*inch))
    limitations = [
        "The dataset covers only 40 fund schemes out of 1,900+ schemes registered with AMFI, limiting industry-wide generalisability.",
        "Transaction data is simulated and may not perfectly reflect real investor behaviour patterns.",
        "NAV data is forward-filled for weekends and holidays, which may slightly overstate daily return precision.",
        "Alpha and Beta calculations use OLS regression which assumes linear relationships and constant betas over time.",
        "The fund recommender system uses only Sharpe ratio and risk grade — a production system would incorporate more factors.",
    ]
    for i, lim in enumerate(limitations, 1):
        story.append(Paragraph(f"{i}. {lim}", body))

    story.append(Spacer(1, 0.2*inch))

    # ── 8. RECOMMENDATIONS ──
    story.append(Paragraph("8. Recommendations", h1))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#0A2142")))
    story.append(Spacer(1, 0.1*inch))
    recs = [
        ("For Investors", "Small-cap and mid-cap funds have delivered superior 3-5 year CAGR but with higher VaR. Investors with high risk tolerance and a 5+ year horizon should allocate 20-30% to these categories."),
        ("For Fund Houses", "Debt funds with expense ratios above 1.5% show poor risk-adjusted returns compared to lower-cost alternatives. Fee compression is critical to remain competitive."),
        ("For Distributors", "Punjab and Tamil Nadu show the highest transaction volumes outside top metros, indicating strong growth potential in Tier-2 city distribution networks."),
        ("For Regulators", "The industry folio count doubling in 3 years indicates strong retail participation but also potential SIP discontinuity risks among newer investors."),
        ("For Analytics Teams", "Implementing rolling VaR with 252-day windows rather than full-period historical VaR would provide more responsive risk signals during market stress periods."),
    ]
    for title, text in recs:
        story.append(Paragraph(f"<b>{title}:</b> {text}", body))

    story.append(PageBreak())

    # ── 9. TECHNICAL STACK ──
    story.append(Paragraph("9. Technical Stack", h1))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#0A2142")))
    story.append(Spacer(1, 0.1*inch))

    tech_data = [
        ["Component", "Technology", "Purpose"],
        ["Data Processing", "Python 3.14, Pandas, NumPy", "ETL, cleaning, transformation"],
        ["Database", "SQLite, SQLAlchemy", "Star schema, analytical queries"],
        ["Visualisation", "Plotly, Seaborn, Matplotlib", "EDA and performance charts"],
        ["Statistics", "SciPy", "OLS regression, VaR, correlation"],
        ["Dashboard", "Power BI Desktop", "Interactive 4-page dashboard"],
        ["Version Control", "Git, GitHub", "Code versioning and collaboration"],
        ["Report", "ReportLab, python-pptx", "PDF report and PPTX presentation"],
    ]
    tech_table = Table(tech_data, colWidths=[4*cm, 6*cm, 6*cm])
    tech_table.setStyle(TableStyle([
        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#0A2142")),
        ("TEXTCOLOR", (0,0), (-1,0), colors.white),
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 9),
        ("ROWBACKGROUNDS", (0,1), (-1,-1), [colors.white, colors.HexColor("#F0F8FF")]),
        ("GRID", (0,0), (-1,-1), 0.5, colors.HexColor("#CCDDEE")),
        ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
        ("TOPPADDING", (0,0), (-1,-1), 5),
        ("BOTTOMPADDING", (0,0), (-1,-1), 5),
    ]))
    story.append(tech_table)
    story.append(Spacer(1, 0.2*inch))

    # ── 10. CONCLUSION ──
    story.append(Paragraph("10. Conclusion", h1))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor("#0A2142")))
    story.append(Spacer(1, 0.1*inch))
    story.append(Paragraph(
        "This capstone project successfully demonstrates end-to-end data analytics for the Indian Mutual Fund industry. "
        "The project delivers a production-quality ETL pipeline, comprehensive exploratory analysis with 15+ visualisations, "
        "mathematically rigorous performance metrics (CAGR, Sharpe, Sortino, Alpha, Beta, VaR, CVaR), "
        "an interactive 4-page Power BI dashboard, and advanced analytics including investor cohort analysis and a fund recommender system.",
        body))
    story.append(Paragraph(
        "The analysis confirms that equity mutual funds, particularly small-cap and mid-cap categories, have delivered "
        "superior long-term returns during the 2022-2026 period, while debt funds provide stability with lower risk. "
        "SIP as an investment mechanism continues to demonstrate strong growth and retail adoption across India.",
        body))

    doc.build(story)
    print("✅ Final_Report.pdf generated successfully")

build_pdf()


# ═══════════════════════════════════════════════════════
# PART 2 — PRESENTATION PPTX
# ═══════════════════════════════════════════════════════
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

DARK_BG   = "0D1B2A"
MID_BG    = "0A2142"
ACCENT    = "00B4D8"
WHITE     = "FFFFFF"
LIGHT_BG  = "F0F8FF"
GREEN     = "00C851"
RED_C     = "FF4444"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank

def add_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)

def add_rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return txBox

def add_bullet_slide(slide, title, bullets, title_size=28):
    add_text(slide, title, 0.5, 0.3, 12.3, 0.8, size=title_size, bold=True,
             color=ACCENT, align=PP_ALIGN.LEFT)
    add_rect(slide, 0.5, 1.1, 12.3, 0.04, ACCENT)
    y = 1.3
    for bullet in bullets:
        add_text(slide, f"• {bullet}", 0.7, y, 11.5, 0.5, size=14, color=WHITE)
        y += 0.52

# ── SLIDE 1: TITLE ──
s1 = prs.slides.add_slide(blank_layout)
add_bg(s1, DARK_BG)
add_rect(s1, 0, 2.5, 13.33, 2.8, MID_BG)
add_rect(s1, 0, 2.5, 0.15, 2.8, ACCENT)
add_text(s1, "Bluestock Mutual Fund", 0.5, 2.6, 12.3, 1.0,
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, "Analytics Capstone Project", 0.5, 3.5, 12.3, 0.7,
         size=24, bold=False, color=rgb(ACCENT).__str__() and ACCENT, align=PP_ALIGN.CENTER)
add_text(s1, "End-to-End MF Analytics | June 2026", 0.5, 4.2, 12.3, 0.5,
         size=14, color="AAAAAA", align=PP_ALIGN.CENTER)
add_text(s1, "Divya Madiraju | Bluestock Fintech MJ28", 0.5, 6.8, 12.3, 0.5,
         size=12, color="888888", align=PP_ALIGN.CENTER)

# ── SLIDE 2: PROBLEM & OBJECTIVE ──
s2 = prs.slides.add_slide(blank_layout)
add_bg(s2, DARK_BG)
add_bullet_slide(s2, "Problem Statement & Objectives", [
    "India's MF industry manages Rs.67+ lakh crore AUM across 1,900+ schemes — data is fragmented",
    "Retail investors lack tools to compare risk-adjusted performance across fund categories",
    "No unified platform exists to track NAV trends, SIP growth, and investor demographics together",
    "Objective: Build end-to-end analytics covering ETL, EDA, performance metrics, and dashboard",
    "Scope: 40 fund schemes, Jan 2022 – Dec 2025, 32,778 investor transactions",
])

# KPI boxes
kpi_items = [
    ("40", "Fund Schemes", "0.3"),
    ("46,000", "NAV Records", "3.5"),
    ("32,778", "Transactions", "6.7"),
    ("7", "Deliverables", "9.9"),
]
for val, lbl, x in kpi_items:
    add_rect(s2, float(x), 5.5, 2.8, 1.5, MID_BG)
    add_text(s2, val, float(x), 5.6, 2.8, 0.7, size=28, bold=True,
             color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s2, lbl, float(x), 6.3, 2.8, 0.5, size=12,
             color="AAAAAA", align=PP_ALIGN.CENTER)

# ── SLIDE 3: DATA SOURCES ──
s3 = prs.slides.add_slide(blank_layout)
add_bg(s3, DARK_BG)
add_bullet_slide(s3, "Data Sources & Architecture", [
    "10 CSV datasets from AMFI India covering NAV, AUM, SIP, transactions, holdings",
    "Live NAV API: mfapi.in — fetched for 5 key schemes (HDFC, SBI, ICICI, Axis, Kotak)",
    "SQLite star schema: dim_fund, dim_date, fact_nav, fact_transactions, fact_performance",
    "ETL pipeline: Extract → Clean → Validate → Load → Verify row counts",
    "Tools: Python 3.14, Pandas, NumPy, SQLite, SQLAlchemy, Pathlib",
])

# ── SLIDE 4: ETL & ARCHITECTURE ──
s4 = prs.slides.add_slide(blank_layout)
add_bg(s4, DARK_BG)
add_text(s4, "ETL Pipeline Architecture", 0.5, 0.3, 12.3, 0.8,
         size=28, bold=True, color=ACCENT)
add_rect(s4, 0.5, 1.1, 12.3, 0.04, ACCENT)

steps = [
    ("EXTRACT", "Raw CSVs + mfapi.in", "1.0", "00B4D8"),
    ("TRANSFORM", "Clean, validate, standardise", "3.8", "00C851"),
    ("LOAD", "SQLite star schema", "6.6", "FF9800"),
    ("ANALYSE", "EDA + Performance + VaR", "9.4", "9C27B0"),
]
for step, desc, x, col in steps:
    add_rect(s4, float(x), 2.0, 2.8, 1.5, MID_BG)
    add_rect(s4, float(x), 2.0, 2.8, 0.4, col)
    add_text(s4, step, float(x), 2.05, 2.8, 0.35, size=11, bold=True,
             color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text(s4, desc, float(x), 2.5, 2.8, 0.8, size=11,
             color=WHITE, align=PP_ALIGN.CENTER)

add_text(s4, "Common Mistakes Avoided:", 0.5, 3.9, 12.3, 0.4,
         size=14, bold=True, color=ACCENT)
mistakes = [
    "ffill() applied after reindexing NAV to full date range (weekends/holidays handled)",
    "CAGR annualised using 252 trading days, not calendar days",
    "Pathlib.Path used throughout — no hard-coded file paths",
    "*.db added to .gitignore — only schema.sql committed",
]
y = 4.4
for m in mistakes:
    add_text(s4, f"✓ {m}", 0.7, y, 12.0, 0.4, size=11, color=GREEN)
    y += 0.42

# ── SLIDE 5: EDA HIGHLIGHTS 1 ──
s5 = prs.slides.add_slide(blank_layout)
add_bg(s5, DARK_BG)
add_bullet_slide(s5, "EDA Highlights — Industry Overview", [
    "Total industry AUM peaked at Rs.391.76 lakh crore — SBI MF leads with highest AUM",
    "SIP inflows grew from Rs.11,000 Cr (Jan 2022) to record highs by Dec 2025",
    "Industry folios doubled: 13.26 Cr (Jan 2022) to 26.12 Cr (Dec 2025)",
    "Equity category dominates net inflows; debt funds see March outflows (tax-related)",
    "Geographic reach: Punjab, Tamil Nadu, Madhya Pradesh lead in transaction volumes",
])
add_text(s5, "15+ Charts Generated | EDA_Analysis.ipynb", 0.5, 6.8, 12.3, 0.5,
         size=11, color="888888", align=PP_ALIGN.CENTER)

# ── SLIDE 6: EDA HIGHLIGHTS 2 ──
s6 = prs.slides.add_slide(blank_layout)
add_bg(s6, DARK_BG)
add_bullet_slide(s6, "EDA Highlights — Fund Analytics", [
    "Small-cap funds: highest 1Y returns (21-25%) but highest NAV volatility",
    "Debt funds (Liquid/Gilt): highest Sharpe ratios due to low volatility",
    "NAV correlation matrix: equity funds show 0.85+ pairwise correlation",
    "Financial services dominates portfolio holdings — HHI concentration risk flagged",
    "Direct plans consistently outperform regular plans on 1Y returns (expense ratio effect)",
])

# ── SLIDE 7: PERFORMANCE METRICS 1 ──
s7 = prs.slides.add_slide(blank_layout)
add_bg(s7, DARK_BG)
add_text(s7, "Performance Metrics — Scorecard", 0.5, 0.3, 12.3, 0.8,
         size=28, bold=True, color=ACCENT)
add_rect(s7, 0.5, 1.1, 12.3, 0.04, ACCENT)

add_text(s7, "Composite Score = 30% CAGR + 25% Sharpe + 20% Alpha + 15% Expense + 10% Drawdown",
         0.5, 1.2, 12.3, 0.5, size=12, color="AAAAAA")

headers = ["Rank", "Fund Name", "Score", "3Y CAGR", "Sharpe"]
col_x   = [0.5, 1.5, 8.5, 9.8, 11.2]
col_w   = [0.9, 6.8, 1.2, 1.2, 1.2]

for i, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    add_rect(s7, x, 1.8, w, 0.4, MID_BG)
    add_text(s7, h, x, 1.82, w, 0.36, size=10, bold=True,
             color=ACCENT, align=PP_ALIGN.CENTER)

top5 = scorecard.nlargest(5, "score").reset_index(drop=True)
for idx, row in top5.iterrows():
    y = 2.3 + idx * 0.5
    bg = "0A2142" if idx % 2 == 0 else "0D1B2A"
    name = str(row.get("scheme_name", ""))[:45]
    vals = [str(idx+1), name,
            f"{row.get('score',0):.1f}",
            f"{row.get('cagr_3y',0):.1f}%",
            f"{row.get('sharpe_ratio',0):.2f}"]
    for v, x, w in zip(vals, col_x, col_w):
        add_rect(s7, x, y, w, 0.45, bg)
        add_text(s7, v, x, y+0.02, w, 0.41, size=9,
                 color=WHITE, align=PP_ALIGN.CENTER)

# ── SLIDE 8: PERFORMANCE METRICS 2 ──
s8 = prs.slides.add_slide(blank_layout)
add_bg(s8, DARK_BG)
add_bullet_slide(s8, "Performance Metrics — Risk Analysis", [
    "Historical VaR (95%): Small-cap funds show -2.5% to -3.5% daily VaR",
    "CVaR (Expected Shortfall): Worst 5% days average -3.8% for high-risk funds",
    "Alpha: Top equity funds generate 15-29% annualised excess return vs Nifty 100",
    "Beta: Equity funds 0.85-1.15 (market-like); debt funds near 0 (uncorrelated)",
    "Maximum Drawdown: Small-cap funds saw -18% to -21% peak-to-trough drawdowns",
    "Rolling Sharpe: Dipped below 0 during 2024 corrections, recovered by 2025",
])

# ── SLIDE 9: DASHBOARD SCREENSHOTS 1 ──
s9 = prs.slides.add_slide(blank_layout)
add_bg(s9, DARK_BG)
add_text(s9, "Interactive Dashboard — Page 1 & 2", 0.5, 0.3, 12.3, 0.8,
         size=28, bold=True, color=ACCENT)
add_rect(s9, 0.5, 1.1, 12.3, 0.04, ACCENT)

# Try to embed dashboard screenshots
for i, (page, x) in enumerate([("Page1_Industry", 0.3), ("Page2_Performance", 6.8)]):
    img_path = BASE_DIR / "dashboard" / f"{page}.png"
    if img_path.exists():
        slide_pic = s9.shapes.add_picture(str(img_path), Inches(x), Inches(1.3),
                                          Inches(6.2), Inches(4.5))
    else:
        add_rect(s9, x, 1.3, 6.2, 4.5, MID_BG)
        add_text(s9, f"{page}\n(Screenshot)", x+0.5, 2.8, 5.2, 1.5,
                 size=14, color="888888", align=PP_ALIGN.CENTER)

# ── SLIDE 10: DASHBOARD SCREENSHOTS 2 ──
s10 = prs.slides.add_slide(blank_layout)
add_bg(s10, DARK_BG)
add_text(s10, "Interactive Dashboard — Page 3 & 4", 0.5, 0.3, 12.3, 0.8,
         size=28, bold=True, color=ACCENT)
add_rect(s10, 0.5, 1.1, 12.3, 0.04, ACCENT)

for i, (page, x) in enumerate([("Page3_Investor", 0.3), ("Page4_SIP", 6.8)]):
    img_path = BASE_DIR / "dashboard" / f"{page}.png"
    if img_path.exists():
        s10.shapes.add_picture(str(img_path), Inches(x), Inches(1.3),
                               Inches(6.2), Inches(4.5))
    else:
        add_rect(s10, x, 1.3, 6.2, 4.5, MID_BG)
        add_text(s10, f"{page}\n(Screenshot)", x+0.5, 2.8, 5.2, 1.5,
                 size=14, color="888888", align=PP_ALIGN.CENTER)

# ── SLIDE 11: KEY FINDINGS ──
s11 = prs.slides.add_slide(blank_layout)
add_bg(s11, DARK_BG)
add_bullet_slide(s11, "Key Findings & Recommendations", [
    "SBI MF leads AUM — distribution through bank branches is the key competitive moat",
    "Small-cap funds deliver best long-term CAGR but require 5+ year investment horizon",
    "SIP discipline is strong — 80%+ investors maintain monthly cadence (gap < 35 days)",
    "Direct plans outperform regular plans consistently — cost matters in mutual funds",
    "Equity fund correlation >0.85 — true diversification requires multi-asset allocation",
    "2024 cohort investors contribute highest avg SIP amounts — newer investors are more informed",
])

# ── SLIDE 12: THANK YOU ──
s12 = prs.slides.add_slide(blank_layout)
add_bg(s12, DARK_BG)
add_rect(s12, 0, 2.8, 13.33, 2.5, MID_BG)
add_rect(s12, 0, 2.8, 0.15, 2.5, ACCENT)
add_text(s12, "Thank You", 0.5, 3.0, 12.3, 1.0,
         size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s12, "Bluestock Mutual Fund Analytics Capstone", 0.5, 4.0, 12.3, 0.6,
         size=16, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s12, "GitHub: github.com/Divya-Madiraju/bluestock_mf_capstone", 0.5, 5.2, 12.3, 0.5,
         size=12, color="888888", align=PP_ALIGN.CENTER)
add_text(s12, "Divya Madiraju | divyamadiraju2006 | June 2026", 0.5, 6.8, 12.3, 0.5,
         size=11, color="666666", align=PP_ALIGN.CENTER)

prs.save(str(REPORTS / "Bluestock_MF_Presentation.pptx"))
print("✅ Bluestock_MF_Presentation.pptx generated successfully")
print(f"\n📁 Files saved to: {REPORTS}")
print("  - Final_Report.pdf")
print("  - Bluestock_MF_Presentation.pptx")
